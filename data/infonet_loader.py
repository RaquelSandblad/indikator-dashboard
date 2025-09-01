"""
Simple loader to extract tables and text from a PowerPoint (pptx) used in the project.
The loader returns a list of tables (as pandas DataFrames) and a concatenated text blob.
"""
from pathlib import Path
from typing import List, Tuple
import pandas as pd

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def load_pptx_tables(filepath: str) -> Tuple[List[pd.DataFrame], str]:
    """Load tables and text from a .pptx file.

    Returns (tables, text_blob). If python-pptx isn't installed or file can't be read,
    returns ([], "").
    """
    p = Path(filepath)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {filepath}")

    if Presentation is None:
        raise RuntimeError("python-pptx not installed. Add python-pptx to requirements.")

    prs = Presentation(str(p))
    tables = []
    text_parts = []

    for slide in prs.slides:
        # Collect slide-level text
        for shape in slide.shapes:
            # Safe text extraction
            try:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        text_parts.append(txt)
            except Exception:
                # Some shapes may raise when accessing text; skip
                pass

            # Safe table extraction
            try:
                tbl = shape.table
            except Exception:
                tbl = None

            if tbl is not None:
                rows = []
                for r in range(len(tbl.rows)):
                    row_cells = []
                    for c in range(len(tbl.columns)):
                        try:
                            cell_text = tbl.cell(r, c).text.strip()
                        except Exception:
                            cell_text = ""
                        row_cells.append(cell_text)
                    rows.append(row_cells)
                if rows:
                    df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows) > 1 else pd.DataFrame(rows)
                    tables.append(df)

    text_blob = "\n\n".join(text_parts)
    return tables, text_blob
